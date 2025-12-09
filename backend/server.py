import os
import io
import uuid
import logging
import re
import difflib
import requests
from datetime import datetime, timezone
from pathlib import Path
from typing import List
import base64

import numpy as np
import cv2
from PIL import Image
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pdf2image import convert_from_bytes

from fastapi import FastAPI, File, UploadFile, APIRouter
from starlette.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field

from motor.motor_asyncio import AsyncIOMotorClient
import easyocr

# ----------------------------
# ENV + Setup
# ----------------------------
ROOT_DIR = Path(__file__).parent

# ----------------------------
# Load environment (.env)
# ----------------------------
if (ROOT_DIR / ".env").exists():
    from dotenv import load_dotenv
    load_dotenv(ROOT_DIR / ".env")

HF_API_KEY = os.getenv("HF_API_KEY")
HF_EMBED_URL = "https://api-inference.huggingface.co/pipeline/feature-extraction/sentence-transformers/all-MiniLM-L6-v2"
HF_TROCR_URL = "https://api-inference.huggingface.co/models/microsoft/trocr-base-printed"
HEADERS = {"Authorization": f"Bearer {HF_API_KEY}"} if HF_API_KEY else {}

if not HF_API_KEY:
    logging.warning("HF_API_KEY not found in environment. Embedding and HF OCR calls will fail.")

# ----------------------------
# Database
# ----------------------------
client = AsyncIOMotorClient(os.getenv("MONGO_URL")) if os.getenv("MONGO_URL") else None
db = client[os.getenv("DB_NAME", "documents")] if client else None

# ----------------------------
# FastAPI App
# ----------------------------
app = FastAPI()
api = APIRouter(prefix="/api")

# ----------------------------
# CORS FIX (IMPORTANT)
# ----------------------------
origins = [
    "https://docusortai-frontend.onrender.com",
    "http://localhost:3000",
    "http://localhost:5173",
    "http://127.0.0.1:3000",
    "http://127.0.0.1:5173",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ----------------------------
# Your API routes go below
# ----------------------------
@app.get("/")
def root():
    return {"status": "backend running"}



def load_reader():
    return easyocr.Reader(['en'], gpu=False)


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ----------------------------
# Document Keywords
# ----------------------------
DOCUMENT_PATTERNS = {
    "Resume": {"keywords": ["experience", "education", "skills", "projects", "work history",
                             "employment", "certifications", "curriculum vitae",
                             "b.tech", "internship", "summary", "achievements"]},
    "Invoice": {"keywords": ["invoice", "invoice no", "bill to", "amount due", "total",
                             "subtotal", "payment", "gst", "vat", "price", "quantity",
                             "description", "receipt"]},
    "Contract": {"keywords": ["contract", "agreement", "party a", "party b", "terms",
                              "conditions", "confidentiality", "governing law",
                              "witness", "signatures", "obligations", "liability",
                              "hereby", "hereto"]},
    "Report": {"keywords": ["report", "site report", "project details", "project report",
                            "summary", "description", "details", "findings", "assessment",
                            "audit", "inspections", "issues", "priority", "status",
                            "comments", "assigned to", "due date", "in progress", "progress"]}
}

# ----------------------------
# Models
# ----------------------------
class ClassificationResult(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    file_name: str
    file_type: str
    file_size: int
    document_type: str
    summary: str
    reason: str
    confidence: float
    extracted_text: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


# ----------------------------
# OCR Helpers
# ----------------------------
def preprocess_image(img):
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.convertScaleAbs(gray, alpha=1.4, beta=10)
    blur = cv2.GaussianBlur(gray, (3, 3), 0)
    _, thresh = cv2.threshold(blur, 0, 255, cv2.THRESH_OTSU + cv2.THRESH_BINARY)
    return thresh


def hf_ocr_fallback(image_bytes: bytes) -> str:
    """Use HF TrOCR as fallback OCR when EasyOCR fails (helps with image invoices)."""
    if not HF_API_KEY:
        return ""
    try:
        encoded = base64.b64encode(image_bytes).decode()
        payload = {"inputs": {"image": encoded}}
        resp = requests.post(HF_TROCR_URL, headers=HEADERS, json=payload, timeout=40)
        data = resp.json()
        if isinstance(data, list) and len(data) > 0 and "generated_text" in data[0]:
            return data[0]["generated_text"] or ""
        # Some models return dict with 'text' key
        if isinstance(data, dict) and "text" in data:
            return data["text"] or ""
        return ""
    except Exception as e:
        logger.warning("HF OCR fallback error: %s", e)
        return ""


def extract_text_from_image(raw):
    """Primary OCR: EasyOCR. If result empty or low words, fallback to HF OCR."""
    try:
        image = Image.open(io.BytesIO(raw)).convert("RGB")
        arr = np.array(image)
        processed = preprocess_image(arr)
        try:
            reader = load_reader() 
            text_list = reader.readtext(processed, detail=0)
        except Exception as e:
            logger.warning("EasyOCR processing error: %s", e)
            text_list = []

        if not text_list:
            # try raw arr
            try:
                reader = load_reader() 
                text_list = reader.readtext(arr, detail=0)
            except Exception:
                text_list = []

        text = " ".join(str(t) for t in text_list).strip()
        # fallback to HF OCR if empty or very short (helps scanned invoice images)
        if (not text) or (len(text.split()) < 3):
            hf_text = hf_ocr_fallback(raw)
            if hf_text and len(hf_text.strip()) > len(text):
                text = hf_text.strip()
        return text
    except Exception as e:
        logger.exception("image extraction error: %s", e)
        # final fallback: try HF OCR
        return hf_ocr_fallback(raw)


def extract_text_from_pdf(raw):
    try:
        reader_pdf = PyPDF2.PdfReader(io.BytesIO(raw))
        text = ""
        for p in reader_pdf.pages:
            text += p.extract_text() or ""
        return text.strip()
    except Exception as e:
        logger.warning("pdf text extract error: %s", e)
        return ""


def extract_text_from_pdf_with_ocr(raw):
    """Hybrid: first try text extraction; then EasyOCR per page; then HF OCR fallback per page."""
    try:
        text = extract_text_from_pdf(raw)
        if text and len(text.strip()) > 20:
            return text.strip()

        pages = convert_from_bytes(raw)
        ocr_text = ""
        for page in pages:
            buf = io.BytesIO()
            page.save(buf, format="PNG")
            page_bytes = buf.getvalue()
            # try EasyOCR on page image
            try:
                arr = np.array(page)
                processed = preprocess_image(arr)
                reader = load_reader() 
                page_text = " ".join(str(t) for t in (reader.readtext(processed, detail=0) or []))
                if not page_text or len(page_text.split()) < 3:
                    page_text = " ".join(str(t) for t in (reader.readtext(np.array(page), detail=0) or []))
            except Exception as e:
                logger.debug("EasyOCR page error: %s", e)
                page_text = ""
            # fallback HF OCR for page if empty
            if not page_text or len(page_text.split()) < 3:
                page_text = hf_ocr_fallback(page_bytes)
            ocr_text += (page_text + "\n")
        return ocr_text.strip()
    except Exception as e:
        logger.exception("PDF OCR error: %s", e)
        return ""


def extract_text_from_docx(raw):
    try:
        doc = DocxDocument(io.BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.warning("docx extract error: %s", e)
        return ""


def extract_text_from_pptx(raw):
    try:
        prs = Presentation(io.BytesIO(raw))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                extract_text_from_shape(shape, texts)
        return "\n".join(t.strip() for t in texts if t.strip())
    except Exception as e:
        logger.warning("pptx extract error: %s", e)
        return ""


def extract_text_from_shape(shape, texts):
    if getattr(shape, "has_text_frame", False):
        texts.append(shape.text_frame.text)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            extract_text_from_shape(subshape, texts)


# ----------------------------
# Email Detection (stricter)
# ----------------------------

def is_contract_hard_match(t: str) -> bool:
    """High-precision contract detection: looks for legal phrasing,
       requires at least two strong contract indicators to avoid false positives."""
    contract_indicators = [
        "agreement", "hereby", "witnesseth", "party a", "party b",
        "governing law", "terms and conditions", "confidentiality",
        "whereas", "hereto", "herein", "indemnif", "obligations", "signatures"
    ]
    count = sum(1 for w in contract_indicators if w in t)
    return count >= 2


def is_invoice_hard_match(t: str) -> bool:
    """High-precision invoice detection:
       require the word 'invoice' plus invoice-like structure (currency, 'amount due', invoice number, or billing table headers)."""
    if "invoice" not in t:
        return False

    # invoice number or bill to / amount due / subtotal etc.
    invoice_no = bool(re.search(r"invoice\s*(no|number|#)\s*[:#\-\s]*\w+", t))
    billing_words = any(x in t for x in ["bill to", "amount due", "amount:", "balance due", "subtotal", "total due"])
    table_headers = any(x in t for x in ["qty", "quantity", "unit price", "description", "unit cost", "amount"])
    # currency pattern: $, ₹, USD, INR, GBP, EUR, numbers with currency words
    currency = bool(re.search(r"[\$\£\€\₹]|(?:\bUSD\b|\bINR\b|\bGBP\b|\bEUR\b|\bAUD\b|\bCAD\b)", t, flags=re.I))
    # numeric line like "Total: 1,234.56" or "Amount Due 1234"
    numeric_total = bool(re.search(r"(total|amount|balance)[\s:\-]*[₹\$\£\€]?\s*\d{1,3}(?:[,\d]{0,})", t, flags=re.I))

    # require at least one structural hint in addition to 'invoice'
    return invoice_no or billing_words or table_headers or currency or numeric_total

def is_email(text: str, debug: bool = False) -> bool:
    t = (text or "").lower()

    # headers must be at start of line to be strong evidence
    header_lines = re.findall(r"(?m)^(from:|to:|subject:|cc:|bcc:|date:)", t)
    header_count = len(header_lines)

    emails = re.findall(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", t)
    email_count = len(emails)

    greetings = ["\nhi ", "\nhello ", "\ndear ", "\nhi,", "\nhello,", "\ndear,"]
    signoffs = ["regards", "sincerely", "best regards", "thanks", "thank you"]

    has_greeting = any(g in t for g in greetings)
    has_signoff = any(s in t for s in signoffs)

    paragraphs = [p.strip() for p in t.split("\n\n") if p.strip()]
    long_paragraphs = sum(1 for p in paragraphs if len(p.split()) > 20)

    # compute base score
    score = 0.0
    score += min(header_count, 3) * 1.5     # header lines are strong
    score += min(email_count, 3) * 0.6
    score += 1.0 if has_greeting else 0.0
    score += 1.0 if has_signoff else 0.0
    score += min(long_paragraphs, 2) * 0.6

    # reduce score if document looks like a report (so reports with contact emails won't be classified as email)
    report_indicators = [
        "site report", "project report", "inspection", "findings", "assessment",
        "audit", "site", "project details", "in progress", "work log", "progress"
    ]
    report_matches = sum(1 for w in report_indicators if w in t)
    score -= min(report_matches, 3) * 1.0  # subtract strong penalty for report indicators

    # reduce if many contract indicators (contracts sometimes include emails)
    contract_indicators = ["agreement", "party a", "party b", "hereby", "witnesseth", "signatures"]
    contract_matches = sum(1 for w in contract_indicators if w in t)
    score -= min(contract_matches, 2) * 0.8

    if debug:
        logger.info("EMAIL DEBUG headers=%s emails=%s greeting=%s signoff=%s long_paras=%s report_matches=%s contract_matches=%s score=%.2f",
                    header_count, email_count, has_greeting, has_signoff, long_paragraphs, report_matches, contract_matches, score)

    # require stronger evidence than before
    return score >= 3.0


# ----------------------------
# Embedding via HF API (safe)
# ----------------------------
def embed_text(text: str):
    """Return normalized embedding vector or zeros on failure."""
    if not text:
        return np.zeros((384,), dtype=float)
    if not HF_API_KEY:
        return np.zeros((384,), dtype=float)
    try:
        payload = {"inputs": text[:500]}
        resp = requests.post(HF_EMBED_URL, headers=HEADERS, json=payload, timeout=30)
        resp.raise_for_status()
        arr = np.array(resp.json()[0], dtype=float)
        norm = np.linalg.norm(arr) + 1e-9
        return arr / norm
    except Exception as e:
        logger.warning("HF embed error: %s", e)
        return np.zeros((384,), dtype=float)


# ----------------------------
# Scoring helpers
# ----------------------------
def keyword_score(text, keywords):
    t = text.lower()
    if not keywords:
        return 0.0
    return sum(1 for k in keywords if k in t) / len(keywords)


def fuzzy_keyword_score(text, keywords, threshold=0.8):
    t = text.lower().split()
    if not keywords:
        return 0.0
    score = 0
    for k in keywords:
        if difflib.get_close_matches(k.lower(), t, n=1, cutoff=threshold):
            score += 1
    return score / len(keywords)


def semantic_score(text, label):
    examples = {
        "Resume": "resume cv work experience education skills projects profile",
        "Invoice": "invoice billing payment gst total amount due",
        "Contract": "legal contract agreement terms conditions parties obligations signatures",
        "Report": "project report site inspection summary findings assessment"
    }
    try:
        emb_text = embed_text(text)
        emb_label = embed_text(examples[label])
        return float(np.dot(emb_text, emb_label))  # cosine-like (already normalized)
    except Exception as e:
        logger.debug("semantic_score error: %s", e)
        return 0.0


# ----------------------------
# Classification Core (improved)
# ----------------------------
def classify(text: str):
    text = (text or "").strip()
    t = text.lower()
    if not t or len(t) < 10:
        return "Others", "Very little text found", 0.2

    # 1) Contract structural check (high precision)
    if is_contract_hard_match(t):
        return "Contract", "Contract wording detected (structural rule)", 0.96

    # 2) Invoice structural check (requires invoice + structure)
    if is_invoice_hard_match(t):
        return "Invoice", "Invoice structural pattern detected", 0.98

    # 3) Email detection (stricter). But allow resume/report overrides below.
    if is_email(text):
        # if report indicators strong, prefer report
        report_kw = keyword_score(t, DOCUMENT_PATTERNS["Report"]["keywords"])
        report_extra = sum(1 for w in ["project details", "site report", "inspection", "findings", "assessment"] if w in t)
        if (report_kw + 0.15 * report_extra) >= 0.45:
            return "Report", "Report indicators override email", min(1.0, report_kw + 0.15 * report_extra)
        # if contract indicators are strong, prefer contract
        if is_contract_hard_match(t):
            return "Contract", "Contract indicators override email", 0.95
        return "Email", "Email structure detected", 0.92

    # 4) Resume quick-check (keep your existing resume logic)
    r_kw = keyword_score(t, DOCUMENT_PATTERNS["Resume"]["keywords"])
    r_sem = semantic_score(text, "Resume") * 0.35
    if (r_kw + r_sem) >= 0.60:
        return "Resume", "Resume-like structure detected", min(1.0, r_kw + r_sem)

    # 5) Full weighted scoring for remaining types
    scores = {}
    reasons = {}
    for doc_type, data in DOCUMENT_PATTERNS.items():
        keywords = data.get("keywords", [])
        if doc_type == "Invoice":
            k = fuzzy_keyword_score(t, keywords)
        else:
            k = keyword_score(t, keywords)
        s = semantic_score(text, doc_type)  # 0..1-ish

        kw_part = 0.5 * k
        sem_part = 0.35 * min(max(s, 0.0), 1.0)
        total = kw_part + sem_part

        # structural bonuses (same as before but slightly tuned)
        bonus = 0.0
        if doc_type == "Report":
            report_indicators = ["project details", "site report", "project report",
                                 "status:", "priority:", "due date", "assigned to",
                                 "description:", "summary:", "issues", "in progress",
                                 "client name", "auditor", "inspection", "findings"]
            b = sum(1 for x in report_indicators if x in t)
            bonus = min(0.20, b * 0.04)
        elif doc_type == "Invoice":
            invoice_indicators = ["invoice", "amount due", "bill", "total", "gst", "vat", "invoice no"]
            b = sum(1 for x in invoice_indicators if x in t)
            bonus = min(0.18, b * 0.06)
        elif doc_type == "Contract":
            contract_indicators = ["agreement", "party a", "party b", "hereby", "witnesseth", "signatures"]
            b = sum(1 for x in contract_indicators if x in t)
            bonus = min(0.20, b * 0.06)
        elif doc_type == "Resume":
            resume_indicators = ["education", "experience", "skills", "projects", "certifications", "curriculum vitae"]
            b = sum(1 for x in resume_indicators if x in t)
            bonus = min(0.15, b * 0.03)

        total += bonus
        scores[doc_type] = total
        reasons[doc_type] = f"kw={k:.2f},sem={s:.3f},kw_part={kw_part:.2f},sem_part={sem_part:.2f},bonus={bonus:.2f}"

    logger.debug("Classification scores: %s", reasons)

    best = max(scores, key=lambda x: scores[x])
    best_score = scores[best]

    if best_score < 0.40:
        return "Others", "Low confidence across all document types", best_score

    return best, "Matched linguistic + semantic patterns", min(best_score, 1.0)

# ----------------------------
# API Routes
# ----------------------------
@api.post("/classify")
async def classify_files(files: List[UploadFile] = File(...)):
    results = []
    for f in files:
        raw = await f.read()
        name = f.filename or "unknown"
        ext = name.split(".")[-1].lower() if "." in name else ""
        if ext == "pdf":
            text = extract_text_from_pdf_with_ocr(raw)
        elif ext in ["docx", "doc"]:
            text = extract_text_from_docx(raw)
        elif ext in ["pptx", "ppt"]:
            text = extract_text_from_pptx(raw)
        elif ext in ["jpg", "jpeg", "png", "webp", "tiff", "bmp"]:
            text = extract_text_from_image(raw)
        elif ext == "txt":
            try:
                text = raw.decode("utf-8", errors="ignore")
            except Exception:
                text = ""
        else:
            # try image OCR fallback
            text = extract_text_from_image(raw)

        doc_type, reason, conf = classify(text)
        summary = f"{doc_type} containing {len(text.split())} words"

        result = ClassificationResult(
            file_name=name,
            file_type=ext.upper(),
            file_size=len(raw),
            document_type=doc_type,
            summary=summary,
            reason=reason,
            confidence=round(float(conf), 3),
            extracted_text=text[:5000]
        )

        # store only if db configured
        if db is not None:
            try:
                doc = result.model_dump()
                doc["timestamp"] = doc["timestamp"].isoformat()
                await db.classifications.insert_one(doc)
            except Exception as e:
                logger.warning("DB insert error: %s", e)

        results.append(result)
    return results


@api.get("/history")
async def history():
    if db is None:
        return []
    items = await db.classifications.find({}, {"_id": 0}).sort("timestamp", -1).limit(50).to_list(50)
    for x in items:
        if isinstance(x.get("timestamp"), str):
            x["timestamp"] = datetime.fromisoformat(x["timestamp"])
    return items

app.include_router(api)
@app.on_event("shutdown")
async def shutdown_db_client():
    if client:
        client.close()


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("server:app", host="0.0.0.0", port=port)
